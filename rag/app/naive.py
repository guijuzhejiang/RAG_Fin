#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
import copy
import io
import logging
import os
import tempfile
import time
import traceback
from pptx import Presentation
import pdf2image
from tika import parser
from io import BytesIO
from docx import Document
from timeit import default_timer as timer
import re

from unstructured.partition.pdf import extractable_elements
from unstructured.partition.pdf_image.ocr import get_table_tokens
from unstructured.partition.pdf_image.pdf_image_utils import get_the_last_modification_date_pdf_or_img, \
    pad_element_bboxes
from unstructured.partition.pdf_image.pdfminer_processing import process_file_with_pdfminer, \
    clean_pdfminer_duplicate_image_elements, clean_pdfminer_inner_elements, process_data_with_pdfminer
from unstructured_inference.inference.elements import Rectangle, region_bounding_boxes_are_almost_the_same
from unstructured_inference.inference.layout import PageLayout, DocumentLayout
from unstructured_inference.inference.layoutelement import LayoutElement
from unstructured_inference.models.tables import cells_to_html

# from deepdoc.parser.pdf_parser_unstructured_VLM import PdfParserVLM, PlainParser
from rag.nlp import rag_tokenizer, naive_merge, tokenize_table, tokenize_chunks, find_codec, concat_img, \
    naive_merge_docx, tokenize_chunks_docx
from deepdoc.parser import PdfParser, PlainParser, ExcelParser, DocxParser, HtmlParser, JsonParser, MarkdownParser, TxtParser, \
    PptParser
from rag.settings import cron_logger
from rag.utils import num_tokens_from_string
from PIL import Image
from functools import reduce
from markdown import markdown
from docx.image.exceptions import UnrecognizedImageError
from typing import List, cast
from unstructured.documents.elements import (
    Text,
)
from unstructured_inference.models import tables
from deepdoc.vision import surya_ocr
from deepdoc.parser.utils import ImageParser, clean_text


class Docx(DocxParser):
    def __init__(self, chat_mdl):
        self.chat_mdl = chat_mdl

    def get_picture(self, document, paragraph):
        img = paragraph._element.xpath('.//pic:pic')
        if not img:
            return None
        img = img[0]
        embed = img.xpath('.//a:blip/@r:embed')[0]
        related_part = document.part.related_parts[embed]
        try:
            image_blob = related_part.image.blob
        except UnrecognizedImageError:
            print("Unrecognized image format. Skipping image.")
            return None
        try:
            image = Image.open(BytesIO(image_blob)).convert('RGB')
            return image
        except Exception as e:
            return None

    def __clean(self, line):
        line = re.sub(r"\u3000", " ", line).strip()
        return line

    def remove_special_characters_len(self, text):
        text_len = 0
        if text:
            # 定义你想删除的特殊字符
            pattern = r'[。？！；!?;+.•]+'
            # 使用 re.sub 函数将这些字符替换为空字符串
            cleaned_text = re.sub(pattern, '', text)
            text_len = len(cleaned_text)
        return text_len

    def __call__(self, filename, binary=None, from_page=0, to_page=100000):
        self.doc = Document(
            filename) if not binary else Document(BytesIO(binary))
        pn = 0
        lines = []
        last_image = None
        img_parser = ImageParser(self.chat_mdl)

        # 获取所有页眉和页脚的内容，用于之后跳过它们
        headers_footers_text = []
        for section in self.doc.sections:
            # 页眉
            if section.header:
                for paragraph in section.header.paragraphs:
                    headers_footers_text.append(paragraph.text)
            # 页脚
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    headers_footers_text.append(paragraph.text)

        for p in self.doc.paragraphs:
            if from_page <= pn < to_page:
                if p.text.strip():
                    if p.style and p.style.name == 'Caption':
                        former_image = None
                        if lines and lines[-1][1] and lines[-1][2] != 'Caption':
                            former_image = lines[-1][1].pop()
                        elif last_image:
                            former_image = last_image
                            last_image = None
                        lines.append((self.__clean(p.text), [former_image], p.style.name))
                    else:
                        current_image = self.get_picture(self.doc, p)
                        image_list = [current_image]
                        if last_image:
                            image_list.insert(0, last_image)
                            last_image = None
                        lines.append((self.__clean(p.text), image_list, p.style.name if p.style else ""))
                else:
                    if current_image := self.get_picture(self.doc, p):
                        if lines:
                            lines[-1][1].append(current_image)
                        else:
                            last_image = current_image
            for run in p.runs:
                if 'lastRenderedPageBreak' in run._element.xml:
                    pn += 1
                    continue
                if 'w:br' in run._element.xml and 'type="page"' in run._element.xml:
                    pn += 1

        # 合并图片的 OCR 结果
        new_line = [(line[0], reduce(concat_img, line[1]) if line[1] else None) for line in lines]
        # new_line = [(line[0], None) for line in lines]

        tbls = []
        for tb in self.doc.tables:
            html = "<table>"
            for r in tb.rows:
                html += "<tr>"
                i = 0
                while i < len(r.cells):
                    span = 1
                    c = r.cells[i]
                    for j in range(i + 1, len(r.cells)):
                        if c.text == r.cells[j].text:
                            span += 1
                            i = j
                    i += 1
                    html += f"<td>{c.text}</td>" if span == 1 else f"<td colspan='{span}'>{c.text}</td>"
                html += "</tr>"
            html += "</table>"
            tbls.append(((None, html), ""))
        return new_line, tbls


class Pdf(object):
    """
    基于unstructured的简化PDF解析器
    使用VLM处理图片和表格，大幅简化原有复杂逻辑
    """
    
    def __init__(self, vlm_mdl=None):
        """
        初始化解析器
        
        Args:
            vlm_mdl: VLM模型实例，用于图片和表格描述
        """
        self.parser = PdfParser(vlm_mdl=vlm_mdl)
        self.outlines = []

    def __call__(self, filename, binary=None, from_page=0,
                 to_page=100000, zoomin=3, callback=None):
        """
        简化的PDF解析方法，使用unstructured库进行解析
        
        Args:
            filename: PDF文件路径
            binary: PDF文件字节流
            from_page: 起始页面（保持兼容性，unstructured会自动处理）
            to_page: 结束页面（保持兼容性，unstructured会自动处理）
            zoomin: 缩放比例（保持兼容性）
            callback: 回调函数
            
        Returns:
            tuple: (sections, tbls) 与原格式完全兼容
        """
        if callback:
            callback(msg="use unstructured and VLM start parsing PDF...")
        
        try:
            # 使用简化解析器
            sections, tbls = self.parser(filename if not binary else binary)
            
            # 复制outlines信息
            self.outlines = getattr(self.parser, 'outlines', [])
            
            if callback:
                callback(0.8, f"Parsing completed: {len(sections)} paragraphs, {len(tbls)} tables")
            
            return sections, tbls
            
        except Exception as e:
            if callback:
                callback(-1, f"解析失败: {str(e)}")
            logging.error(f"PDF解析失败: {str(e)}")
            return [("PDF解析失败", "")], []
    
    def crop(self, text, ZM=3, need_position=False):
        """保持兼容性的crop方法"""
        return self.parser.crop(text, ZM, need_position)
    
    def remove_tag(self, txt):
        """保持兼容性的标签移除方法"""
        return self.parser.remove_tag(txt)
    
    @staticmethod
    def total_page_number(filename, binary=None):
        """获取PDF总页数"""
        return PdfParser.total_page_number(filename, binary)



class Markdown(MarkdownParser):
    def __call__(self, filename, binary=None):
        txt = ""
        tbls = []
        if binary:
            encoding = find_codec(binary)
            txt = binary.decode(encoding, errors="ignore")
        else:
            with open(filename, "r") as f:
                txt = f.read()
        remainder, tables = self.extract_tables_and_remainder(f'{txt}\n')
        sections = []
        tbls = []
        for sec in remainder.split("\n"):
            if num_tokens_from_string(sec) > 10 * self.chunk_token_num:
                sections.append((sec[:int(len(sec) / 2)], ""))
                sections.append((sec[int(len(sec) / 2):], ""))
            else:
                sections.append((sec, ""))
        print(tables)
        for table in tables:
            tbls.append(((None, markdown(table, extensions=['markdown.extensions.tables'])), ""))
        return sections, tbls


def chunk(filename, binary=None, from_page=0, to_page=100000,
          lang="Chinese", callback=None, **kwargs):
    """
        Supported file formats are docx, pdf, excel, txt.
        This method apply the naive ways to chunk files.
        Successive text will be sliced into pieces using 'delimiter'.
        Next, these successive pieces are merge into chunks whose token number is no more than 'Max token number'.
    """

    eng = lang.lower() == "english"  # is_english(cks)
    parser_config = kwargs.get(
        "parser_config", {"chunk_token_num": 128, "delimiter": "\n!?。；！？", "layout_recognize": True})
    chat_mdl = kwargs.get('chat_mdl', None)
    doc = {
        "docnm_kwd": filename,
        "title_tks": rag_tokenizer.tokenize(re.sub(r"\.[a-zA-Z]+$", "", filename))
    }
    doc["title_sm_tks"] = rag_tokenizer.fine_grained_tokenize(doc["title_tks"])
    res = []
    pdf_parser = None
    chunk_token_num = parser_config.get("chunk_token_num", 128)
    if re.search(r"\.(doc|docx|ppt|pptx|)$", filename, re.IGNORECASE):
        callback(0.1, f"Start to parse {filename.split('.')[-1]}.")

        with tempfile.TemporaryDirectory() as temp_dir:
            tmp_filename = str(time.time())
            tmp_filepath = os.path.join(str(temp_dir), f"{tmp_filename}.{filename.split('.')[-1]}")
            file_type = filename.split('.')[-1]
            need_convert = False

            if file_type == "doc" or file_type == "ppt":
                need_convert = True
                tmp_doc_filepath = os.path.join(str(temp_dir), f"{tmp_filename}.{file_type}")
                tmp_filepath = os.path.join(str(temp_dir), f"{tmp_filename}.{file_type}x")
                with open(tmp_doc_filepath, "wb") as f:
                    f.write(binary)
                os.system(f"/usr/bin/libreoffice --headless --convert-to {file_type}x {tmp_doc_filepath} --outdir {str(temp_dir)}")

            if 'doc' in file_type:
                doc_obj = Document(tmp_filepath if need_convert else BytesIO(binary))
                for section in doc_obj.sections:
                    section.different_first_page_header_footer = False
                    section.header.is_linked_to_previous = True
                    # 页眉
                    if section.header:
                        for paragraph in section.header.paragraphs:
                            paragraph.text = ""
                    # 页脚
                    if section.footer:
                        for paragraph in section.footer.paragraphs:
                            paragraph.text = ""
                doc_obj.save(tmp_filepath)
            else:
                """去除 PPT 文件中的页脚和页码"""
                prs = Presentation(tmp_filepath if need_convert else BytesIO(binary))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        # Check if the shape is a placeholder and if it's for footer
                        if shape.is_placeholder and shape.placeholder_format.idx in (
                        1, 2):  # 1 for Footer, 2 for Slide Number
                            sp_text = shape.text_frame.text
                            # print(f"Removing footer: '{sp_text}'")
                            # Clear the text
                            shape.text_frame.clear()
                prs.save(tmp_filepath)
            # with open(tmp_filepath, 'wb') as tmp_f:
            #     tmp_f.write(binary)
            os.system(f"/usr/bin/libreoffice --headless --convert-to pdf {tmp_filepath} --outdir {str(temp_dir)}")
            with open(os.path.join(temp_dir, f"{tmp_filename}.pdf"), 'rb') as tmp_pdf_f:
                pdf_parser = Pdf(
                ) if parser_config.get("layout_recognize", True) else PlainParser()
                sections, tbls = pdf_parser(tmp_pdf_f.read(),
                                            from_page=from_page, to_page=to_page, callback=callback)

    # if re.search(r"\.docx$", filename, re.IGNORECASE):
    #     callback(0.1, "Start to parse Docx.")
    #     sections, tbls = Docx(chat_mdl)(filename, binary)
    #     res = tokenize_table(tbls, doc, eng)  # just for table
    #
    #     callback(0.8, "Finish parsing.")
    #     st = timer()
    #     chunks, images = naive_merge_docx(
    #         sections, int(parser_config.get(
    #             "chunk_token_num", 128)), parser_config.get(
    #             "delimiter", "\n!?。；！？"))
    #
    #     if kwargs.get("section_only", False):
    #         return chunks
    #
    #     res.extend(tokenize_chunks_docx(chunks, doc, eng, images))
    #     cron_logger.info("naive_merge({}): {}".format(filename, timer() - st))
    #     return res

    # elif re.search(r"\.ppt(?:x)?$", filename, re.IGNORECASE):
    #     callback(0.1, "Start to parse ppt.")
    #     sections, tbls = PptParser(chat_mdl)(filename, binary, from_page=0, to_page=1000)
    #     res = tokenize_table(tbls, doc, eng)  # just for table
    #     callback(0.8, "Finish parsing.")

    elif re.search(r"\.pdf$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse pdf.")
        # 传递VLM模型给PDF解析器
        vlm_mdl = kwargs.get('vlm_mdl', None)
        pdf_parser = Pdf(vlm_mdl=vlm_mdl) if parser_config.get("layout_recognize", True) else PlainParser()
        sections, tbls = pdf_parser(filename if not binary else binary,
                                    from_page=from_page, to_page=to_page, callback=callback)
        res = tokenize_table(tbls, doc, eng)
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.xlsx$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse xlsx.")
        excel_parser = ExcelParser()
        sections = [(l, "") for l in excel_parser.html(binary, chat_mdl=chat_mdl, chunk_rows=chunk_token_num) if l]
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.xls$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse xls.")
        excel_parser = ExcelParser()
        sections = [(l, "") for l in
                    excel_parser.html(binary, chat_mdl=chat_mdl, chunk_rows=chunk_token_num, type='xls') if l]
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.(txt|py|js|java|c|cpp|h|php|go|ts|sh|cs|kt|sql)$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections = TxtParser()(filename, binary,
                               parser_config.get("chunk_token_num", 128),
                               parser_config.get("delimiter", "\n!?;。；！？"))
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.(md|markdown)$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections, tbls = Markdown(int(parser_config.get("chunk_token_num", 128)))(filename, binary)
        res = tokenize_table(tbls, doc, eng)
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.(htm|html)$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse html.")
        sections = HtmlParser(chat_mdl)(filename, binary)
        sections = [(l, "") for l in sections if l]
        callback(0.8, "Finish parsing.")

    elif re.search(r"\.json$", filename, re.IGNORECASE):
        callback(0.1, "Start to parse.")
        sections = JsonParser(int(parser_config.get("chunk_token_num", 128)))(binary)
        # sections = [(_, "") for _ in sections if _]
        callback(0.8, "Finish parsing.")

    # elif re.search(r"\.doc$", filename, re.IGNORECASE):
    #     callback(0.1, "Start to parse doc.")
    #     binary = BytesIO(binary)
    #     doc_parsed = parser.from_buffer(binary)
    #     sections = doc_parsed['content'].split('\n')
    #     sections = [(l, "") for l in sections if l]
    #     callback(0.8, "Finish parsing.")

    else:
        raise NotImplementedError(
            "file type not supported yet(pdf, xlsx, doc, docx, txt supported)")

    st = timer()
    if re.search(r"\.json$", filename, re.IGNORECASE):
        chunks = sections
    else:
        chunks = naive_merge(
            sections, int(parser_config.get("chunk_token_num", 128)), parser_config.get("delimiter", "\n!?。；！？"))
    if kwargs.get("section_only", False):
        return chunks

    res.extend(tokenize_chunks(chunks, doc, eng, pdf_parser))
    cron_logger.info("naive_merge({}): {}".format(filename, timer() - st))
    return res


if __name__ == "__main__":
    import sys

    def dummy(prog=None, msg=""):
        pass

    chunk(sys.argv[1], from_page=0, to_page=10, callback=dummy)
