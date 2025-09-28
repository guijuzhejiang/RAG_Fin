#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
import subprocess
from io import BytesIO
from pptx import Presentation
from PIL import Image
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE
from deepdoc.parser.utils import ImageParser


class RAGFlowPptParser(object):
    def __init__(self, chat_mdl):
        super().__init__()
        self.chat_mdl = chat_mdl

    def __extract_table_as_html(self, table):
        """
        Convert the given PPT table into an HTML table format, handling merged cells with both colspan and rowspan.
        """
        html = "<table border='1' cellspacing='0' cellpadding='5'>"

        row_index = 0
        for row in table.rows:
            html += "<tr>"
            cell_index = 0  # 手动维护当前单元格的索引
            while cell_index < len(row.cells):
                cell = row.cells[cell_index]

                # 如果单元格是合并的一部分且不是起始单元格，跳过它
                if cell.is_spanned:
                    cell_index += 1
                    continue

                # 处理 colspan 和 rowspan
                colspan = 1
                rowspan = 1

                # 如果是合并起始单元格，检查横向和纵向合并
                if cell.is_merge_origin:
                    # 计算横向合并（colspan）
                    for next_cell_index in range(cell_index + 1, len(row.cells)):
                        next_cell = row.cells[next_cell_index]
                        if next_cell.is_spanned:
                            colspan += 1
                        else:
                            break

                    # 计算纵向合并（rowspan）
                    for next_row_index in range(row_index + 1, len(table.rows)):
                        next_row_cell = table.rows[next_row_index].cells[cell_index]
                        if next_row_cell.is_spanned:
                            rowspan += 1
                        else:
                            break

                # 生成 HTML 代码
                if colspan > 1 or rowspan > 1:
                    html += f"<td colspan='{colspan}' rowspan='{rowspan}'>{cell.text.strip()}</td>"
                else:
                    html += f"<td>{cell.text.strip()}</td>"

                cell_index += 1  # 移动到下一个单元格

            html += "</tr>"
            row_index += 1  # 移动到下一行

        html += "</table>"
        return html

    def __clean(self, line):
        """
        Clean up the text by removing unwanted characters.
        """
        line = re.sub(r"\u3000", " ", line).strip()
        line = re.sub(r"\n", " ", line).strip()
        return line

    def __call__(self, filename, binary=None, from_page=0, to_page=100000):
        """
        Main method to parse PPT slides and extract text and tables.
        """
        # Load the PPT file
        ppt = Presentation(filename) if not binary else Presentation(BytesIO(binary))
        pn = 0
        lines = []
        tbls = []
        pic_num=0

        # Extract slides content and images
        for i, slide in enumerate(ppt.slides):
            if pn > to_page:
                break
            if from_page <= pn < to_page:
                texts = []
                for shape in sorted(slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left)):
                    # Handle table shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE or shape.has_table:  # 19
                        tbl_html = self.__extract_table_as_html(shape.table)
                        tbls.append(((None, tbl_html), ""))
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE: #13
                        pic_num +=1
                        image_blob = shape.image.blob
                        img_parser = ImageParser(self.chat_mdl)
                        ocr_text = img_parser.get_picture(image_blob)
                        if ocr_text:
                            texts.append(ocr_text)
                    # Handle text and images
                    elif shape.has_text_frame and shape.text_frame.text.strip():
                            texts.append(self.__clean(shape.text_frame.text))
                lines.extend(texts)
            # Move to the next slide
            pn += 1
        print(f"{pic_num} Picture was found !!")
        return lines, tbls
